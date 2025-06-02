# file: backend/core/extractors/pptx_extractor.py

from core.extractors.base_extractor import BaseExtractor
from pptx import Presentation

class PPTXExtractor(BaseExtractor):
    def extract_text(self) -> str:
        try:
            self.file_stream.seek(0)
            prs = Presentation(self.file_stream)
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        txt = shape.text.strip()
                        if txt:
                            text_parts.append(txt)
            return " ".join(" ".join(text_parts).split())
        except Exception as e:
            print(f"❌ PPTXExtractor HATA: {e}")
            return ""
